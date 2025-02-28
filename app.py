import os
import re
import shutil
import logging
from pathlib import Path
from typing import Dict, List, Optional, Any
from contextlib import contextmanager

from fastapi import FastAPI, UploadFile, File, HTTPException, Depends, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field
import uvicorn
import pdfplumber
from docx import Document  # For DOCX text extraction
from pptx import Presentation  # For PPTX text extraction

from langchain_community.llms import Ollama
from langchain_community.vectorstores import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.chains import create_retrieval_chain
from langchain.prompts import PromptTemplate
from langchain.embeddings import HuggingFaceEmbeddings
# Import the LangChain Document schema (aliasing to avoid conflict with python-docx Document)
from langchain.schema import Document as LC_Document

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

# Application configuration
class Config:
    FOLDER_PATH = Path("db")
    UPLOAD_DIR = Path("temp_uploads")
    MODEL_NAME = "llama3.1:latest"
    EMBEDDING_MODEL = "all-MiniLM-L6-v2"
    CHUNK_SIZE = 1024
    CHUNK_OVERLAP = 80
    RETRIEVER_K = 5

    # Ensure directories exist
    @classmethod
    def init_directories(cls):
        cls.FOLDER_PATH.mkdir(exist_ok=True)
        cls.UPLOAD_DIR.mkdir(exist_ok=True)

# Initialize configuration
Config.init_directories()

# Pydantic models
class QueryRequest(BaseModel):
    query: str = Field(..., description="The query text to process")

class Source(BaseModel):
    source: str
    page_content: str

class QueryResponse(BaseModel):
    answer: str
    sources: Optional[List[Source]] = None

class FileResponse(BaseModel):
    filename: str
    extracted_data: Dict[str, str]

# Initialize FastAPI app
app = FastAPI(title="Document Processing API")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Adjust for production security
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Dependency for LLM
@contextmanager
def get_llm():
    try:
        llm = Ollama(model=Config.MODEL_NAME)
        yield llm
    except Exception as e:
        logger.error(f"Error creating LLM: {e}")
        raise HTTPException(status_code=500, detail="Failed to initialize LLM")

# Dependency for embeddings
@contextmanager
def get_embeddings():
    try:
        embeddings = HuggingFaceEmbeddings(model_name=Config.EMBEDDING_MODEL)
        yield embeddings
    except Exception as e:
        # Check if it's a database schema error
        if "no such column: collections.topic" in str(e):
            logger.warning("Database schema error in embeddings. Resetting vector store...")
            if Config.FOLDER_PATH.exists():
                shutil.rmtree(Config.FOLDER_PATH)
                Config.FOLDER_PATH.mkdir(exist_ok=True)
            # Try again with a clean database
            try:
                embeddings = HuggingFaceEmbeddings(model_name=Config.EMBEDDING_MODEL)
                yield embeddings
                return
            except Exception as e2:
                logger.error(f"Error creating embeddings after DB reset: {e2}")
        logger.error(f"Error creating embeddings model: {e}")
        raise HTTPException(status_code=500, detail="Failed to initialize embeddings model")

# Dependency for text splitter
def get_text_splitter():
    return RecursiveCharacterTextSplitter(
        chunk_size=Config.CHUNK_SIZE,
        chunk_overlap=Config.CHUNK_OVERLAP,
        length_function=len,
        is_separator_regex=False
    )

# Dependency for prompt template
def get_prompt_template():
    return PromptTemplate.from_template(
        """
        <s>[INST] You are a helpful assistant tasked with answering questions based on the provided context.
        Please follow these guidelines:
        - Answer only based on the information provided in the context
        - If the information is not in the context, say "I don't have enough information to answer this question"
        - Be concise and direct in your answers
        - Do not make assumptions beyond what is stated in the context
        Question: {input}
        Context: {context}
        Answer: [/INST]</s>
        """
    )

@contextmanager
def get_vector_store():
    try:
        with get_embeddings() as embedding_fn:
            if Config.FOLDER_PATH.exists() and any(Config.FOLDER_PATH.iterdir()):
                try:
                    vector_store = Chroma(
                        persist_directory=str(Config.FOLDER_PATH),
                        embedding_function=embedding_fn
                    )
                    yield vector_store
                except Exception as e:
                    # If there's a database schema error, try to reset the database
                    if "no such column: collections.topic" in str(e):
                        logger.warning("Database schema error detected. Resetting vector store...")
                        if Config.FOLDER_PATH.exists():
                            shutil.rmtree(Config.FOLDER_PATH)
                            Config.FOLDER_PATH.mkdir(exist_ok=True)
                        yield None
                    else:
                        raise
            else:
                yield None
    except Exception as e:
        logger.error(f"Error accessing vector store: {e}")
        raise HTTPException(status_code=500, detail="Failed to access vector store")

# Utility Functions
def clean_text(text: str) -> str:
    """Clean extracted text by removing excessive whitespace."""
    if not text:
        return ""
    return re.sub(r'\s+', ' ', text).strip()

def extract_text_from_pdf(pdf_path: Path) -> str:
    """Extract text from a PDF using pdfplumber."""
    extracted_text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                extracted_text += clean_text(page_text) + " "
        return extracted_text
    except Exception as e:
        logger.error(f"Error processing PDF {pdf_path}: {e}")
        return extracted_text  # Return whatever text was extracted before error

def extract_text_from_word(docx_path: Path) -> str:
    """Extract text from a Word document."""
    extracted_text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            extracted_text += clean_text(para.text) + " "
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    extracted_text += clean_text(cell.text) + " "
        return extracted_text
    except Exception as e:
        logger.error(f"Error processing DOCX {docx_path}: {e}")
        return extracted_text

def extract_text_from_ppt(ppt_path: Path) -> str:
    """Extract text from a PowerPoint presentation."""
    extracted_text = ""
    try:
        pres = Presentation(ppt_path)
        for slide in pres.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(clean_text(shape.text))
            if slide_text:
                extracted_text += " ".join(slide_text) + " "
        return extracted_text
    except Exception as e:
        logger.error(f"Error processing PPTX {ppt_path}: {e}")
        return extracted_text

def process_documents(input_directory: Path) -> Dict[str, str]:
    """Process all supported documents in a directory and extract text."""
    extracted_data = {}
    for file_path in input_directory.glob('**/*'):
        if not file_path.is_file():
            continue

        file_name = file_path.stem
        file_ext = file_path.suffix.lower()

        if file_ext == '.pdf':
            text = extract_text_from_pdf(file_path)
            extracted_data[file_name] = text
        elif file_ext == '.docx':
            text = extract_text_from_word(file_path)
            extracted_data[file_name] = text
        elif file_ext == '.pptx':
            text = extract_text_from_ppt(file_path)
            extracted_data[file_name] = text

    return extracted_data

def create_document_chunks(extracted_data: Dict[str, str], text_splitter) -> List[LC_Document]:
    """Split extracted text into chunks for embedding."""
    docs = []
    for filename, text in extracted_data.items():
        if not text:
            logger.warning(f"No text extracted from {filename}")
            continue

        chunks = text_splitter.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(LC_Document(
                page_content=chunk,
                metadata={"source": filename, "chunk_id": i}
            ))
    return docs

async def save_upload_file(file: UploadFile) -> Path:
    """Save an uploaded file to disk."""
    file_path = Config.UPLOAD_DIR / file.filename
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)
    return file_path

async def process_and_embed_documents(background_tasks: BackgroundTasks, file_path: Path):
    """Process documents and add to vector store (background task)."""
    try:
        # Process documents
        extracted_data = process_documents(Config.UPLOAD_DIR)

        # Skip if no text was extracted
        if not extracted_data:
            logger.warning(f"No text extracted from {file_path}")
            return

        # Create document chunks
        text_splitter = get_text_splitter()
        docs = create_document_chunks(extracted_data, text_splitter)

        if not docs:
            logger.warning("No document chunks created")
            return

        # Add to vector store
        with get_embeddings() as embedding_fn:
            if Config.FOLDER_PATH.exists() and any(Config.FOLDER_PATH.iterdir()):
                vector_store = Chroma(
                    persist_directory=str(Config.FOLDER_PATH),
                    embedding_function=embedding_fn
                )
                vector_store.add_documents(docs)
                vector_store.persist()
            else:
                Chroma.from_documents(
                    docs,
                    embedding_fn,
                    persist_directory=str(Config.FOLDER_PATH)
                )

        logger.info(f"Successfully processed and embedded {file_path}")
    except Exception as e:
        logger.error(f"Error in background processing task: {e}")

# API Endpoints
@app.post("/ai", response_model=QueryResponse)
async def ai_query(query_req: QueryRequest):
    """Query the LLM directly without using retrieved documents."""
    logger.info(f"Processing direct query: {query_req.query}")

    try:
        with get_llm() as llm:
            response = llm.invoke(query_req.query)
            return {"answer": response}
    except Exception as e:
        logger.error(f"Error processing query: {e}")
        raise HTTPException(status_code=500, detail="Failed to process query")

@app.post("/ask_content", response_model=QueryResponse)
async def ask_content(query_req: QueryRequest):
    """Query the LLM using document retrieval for context."""
    logger.info(f"Processing retrieval query: {query_req.query}")

    try:
        with get_vector_store() as vector_store:
            if not vector_store:
                raise HTTPException(
                    status_code=400,
                    detail="No documents have been uploaded. Please upload documents first."
                )

            retriever = vector_store.as_retriever(
                search_type="similarity",
                search_kwargs={"k": Config.RETRIEVER_K},
            )

            with get_llm() as llm:
                prompt = get_prompt_template()
                document_chain = create_stuff_documents_chain(llm, prompt)
                chain = create_retrieval_chain(retriever, document_chain)

                result = chain.invoke({"input": query_req.query})

                # Format sources for response
                sources = []
                for doc in result.get("context", []):
                    sources.append({
                        "source": doc.metadata.get("source", "Unknown"),
                        "page_content": doc.page_content
                    })

                return {"answer": result.get("answer", ""), "sources": sources}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error processing retrieval query: {e}")
        raise HTTPException(status_code=500, detail="Failed to process query")

@app.post("/pdf", response_model=FileResponse)
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...)
):
    """Upload a document (PDF, DOCX, PPTX) for processing and embedding."""
    logger.info(f"Processing uploaded file: {file.filename}")

    # Validate file extension
    if not file.filename.lower().endswith(('.pdf', '.docx', '.pptx')):
        raise HTTPException(
            status_code=400,
            detail="Only PDF, DOCX, and PPTX files are supported"
        )

    try:
        # Save uploaded file
        file_path = await save_upload_file(file)

        # Process documents synchronously to get extracted text
        extracted_data = process_documents(Config.UPLOAD_DIR)

        # Schedule embedding as a background task
        background_tasks.add_task(process_and_embed_documents, background_tasks, file_path)

        return {
            "filename": file.filename,
            "extracted_data": extracted_data
        }
    except Exception as e:
        logger.error(f"Error processing uploaded file: {e}")
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")

@app.delete("/delete_embeddings")
async def delete_embeddings():
    """Delete all stored embeddings."""
    logger.info("Deleting all embeddings")

    try:
        if Config.FOLDER_PATH.exists():
            shutil.rmtree(Config.FOLDER_PATH)
            Config.FOLDER_PATH.mkdir(exist_ok=True)
            return {"status": "Successfully deleted all embeddings"}
        else:
            return {"status": "No embeddings to delete"}
    except Exception as e:
        logger.error(f"Error deleting embeddings: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to delete embeddings: {str(e)}")

@app.on_event("shutdown")
def shutdown_event():
    import torch
    if hasattr(torch, "mps") and torch.mps.is_available():
        # Try to clean up MPS resources
        torch.mps.empty_cache()

if __name__ == "__main__":
    uvicorn.run("app:app", host="0.0.0.0", port=8081, reload=True)
